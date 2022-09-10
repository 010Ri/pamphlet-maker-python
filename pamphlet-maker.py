import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Mm
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import datetime
import glob
import random

# ------------------------------------------------------------------------------------------------
# 白紙のページを追加した後、レイアウト枠の画像を挿入するための関数
def add_page(layout_flame):
    # 白紙のページを追加
    slide_layout = pt.slide_layouts[6]
    slide = pt.slides.add_slide(slide_layout)
    # 画像の挿入
    file_path = "./img/layout-flame/" + layout_flame + ".png"
    pic = slide.shapes.add_picture(file_path, 0, 0, Cm(18.2), None)  # (file path, x座標, y座標, 横幅, 縦幅)
    # add_pictureの引数について：widthかheightのどちらかさえ指定すれば、元のアスペクト比（縦横比）を自動的に保ってくれるらしい
    # 画像の回転
    # pic.rotation = 20

# ------------------------------------------------------------------------------------------------
# layout-flame_1-L , layout-flame_1(8)-Lのテキストボックス調整用の関数
def insert_textbox_1_L(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'  # テキストボックスにあらかじめテキストを入力できる
    textbox.rotation = -90  # テキストボックスを回転させる
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # テキストを左右中央揃えにする
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # テキストを上下中央揃えにする
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE 
    textbox.top = Cm(12.35)
    textbox.left = Cm(-2.73)

# layout-flame_1-R , layout-flame_1(8)-Rのテキストボックス調整用の関数
def insert_textbox_1_R(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(12.35)
    textbox.left = Cm(10.95)

# layout-flame_2-2-Lのテキストボックス調整用の関数
def insert_textbox_2_2_L(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入（上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(6.535)
    textbox.left = Cm(-2.7)
    # テキストボックスの挿入（下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(18.195)
    textbox.left = Cm(-2.7)

# layout-flame_2-2-Rのテキストボックス調整用の関数
def insert_textbox_2_2_R(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入（上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(6.535)
    textbox.left = Cm(10.95)
    # テキストボックスの挿入（下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(18.195)
    textbox.left = Cm(10.95)

# layout-flame_2-4-4-L , layout-flame_2-4-8-Lのテキストボックス調整用の関数
def insert_textbox_2_4_L(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入（上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(6.535)
    textbox.left = Cm(-2.65)
    # テキストボックスの挿入（中側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(15.3)
    textbox.left = Cm(-0.165)
    # テキストボックスの挿入（下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(21.1)
    textbox.left = Cm(-0.165)

# layout-flame_2-4-4-R , layout-flame_2-4-8-Rのテキストボックス調整用の関数
def insert_textbox_2_4_R(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入（上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(6.535)
    textbox.left = Cm(10.95)
    # テキストボックスの挿入（中側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(15.12)
    textbox.left = Cm(13.45)
    # テキストボックスの挿入（下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(20.92)
    textbox.left = Cm(13.45)

# layout-flame_3-4-L , layout-flame_3-8-L , layout-flame_3(4)-4-L , layout-flame_3(4)-8-L , layout-flame_3(8)-4-Lのテキストボックス調整用の関数
def insert_textbox_3_L(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入（上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(9.55)
    textbox.left = Cm(-2.7)
    # テキストボックスの挿入（下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(21.1)
    textbox.left = Cm(-0.165)

# layout-flame_3-4-R , layout-flame_3-8-R , layout-flame_3(4)-4-R , layout-flame_3(4)-8-R , layout-flame_3(8)-8-Rのテキストボックス調整用の関数
def insert_textbox_3_R(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入（上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(10), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告(会社名)・イベント名'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(9.4)
    textbox.left = Cm(10.95)
    # テキストボックスの挿入（下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(21)
    textbox.left = Cm(13.45)

# layout-flame_4-4-4-4-L , layout-flame_4-4-4-8-Lのテキストボックス調整用の関数
def insert_textbox_4_L(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入（上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(3.8)
    textbox.left = Cm(0.1)
    # テキストボックスの挿入（中上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(9.567)
    textbox.left = Cm(0.1)
    # テキストボックスの挿入（中下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(15.433)
    textbox.left = Cm(0.1)
    # テキストボックスの挿入（下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = -90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(21.2)
    textbox.left = Cm(0.1)

# layout-flame_4-4-4-4-R , layout-flame_4-4-4-8-Rのテキストボックス調整用の関数
def insert_textbox_4_R(id):
    slide = pt.slides[id]  # layout-flame_something
    # テキストボックスの挿入（上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(3.7)
    textbox.left = Cm(13.08)
    # テキストボックスの挿入（中上側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(9.52)
    textbox.left = Cm(13.08)
    # テキストボックスの挿入（中下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(15.25)
    textbox.left = Cm(13.08)
    # テキストボックスの挿入（下側）
    textbox = slide.shapes.add_textbox(0, 0, Cm(5), Cm(1))  # (x座標, y座標, 横幅, 縦幅)
    textbox.text = '広告'
    textbox.rotation = 90
    text_frame = textbox.text_frame
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.autosize = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    textbox.top = Cm(21.05)
    textbox.left = Cm(13.08)

# ------------------------------------------------------------------------------------------------
# layout-flame_1-Lの画像調整用の関数
def insert_ad_1_L(id, ad):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(1.54), Cm(13.61), Cm(22.63))

# layout-flame_1-Rの画像調整用の関数
def insert_ad_1_R(id, ad):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad
    slide.shapes.add_picture(file_path, Cm(1.62), Cm(1.54), Cm(13.61), Cm(22.63))

# layout-flame_1(8)-Lの画像調整用の関数
def insert_ad_18_L(id, ad1, ad2, ad3, ad4, ad5, ad6, ad7, ad8):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(1.54), Cm(6.63), Cm(5.58))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(7.32), Cm(6.63), Cm(5.58))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(12.98), Cm(6.63), Cm(5.58))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(18.56), Cm(6.63), Cm(5.58))
    file_path = ad5
    slide.shapes.add_picture(file_path, Cm(9.91), Cm(1.54), Cm(6.63), Cm(5.58))
    file_path = ad6
    slide.shapes.add_picture(file_path, Cm(9.91), Cm(7.32), Cm(6.63), Cm(5.58))
    file_path = ad7
    slide.shapes.add_picture(file_path, Cm(9.91), Cm(12.98), Cm(6.63), Cm(5.58))
    file_path = ad8
    slide.shapes.add_picture(file_path, Cm(9.91), Cm(18.56), Cm(6.63), Cm(5.58))

# layout-flame_1(8)-Rの画像調整用の関数
def insert_ad_18_R(id, ad1, ad2, ad3, ad4, ad5, ad6, ad7, ad8):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(1.54), Cm(6.63), Cm(5.58))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(7.22), Cm(6.63), Cm(5.58))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(12.85), Cm(6.63), Cm(5.58))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(18.59), Cm(6.63), Cm(5.58))
    file_path = ad5
    slide.shapes.add_picture(file_path, Cm(8.49), Cm(1.54), Cm(6.63), Cm(5.58))
    file_path = ad6
    slide.shapes.add_picture(file_path, Cm(8.49), Cm(7.22), Cm(6.63), Cm(5.58))
    file_path = ad7
    slide.shapes.add_picture(file_path, Cm(8.49), Cm(12.85), Cm(6.63), Cm(5.58))
    file_path = ad8
    slide.shapes.add_picture(file_path, Cm(8.49), Cm(18.59), Cm(6.63), Cm(5.58))

# layout-flame_2-2-Lの画像調整用の関数
def insert_ad_2_2_L(id, ad1, ad2):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(1.54), Cm(13.56), Cm(11.01))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(13.15), Cm(13.56), Cm(11.01))

# layout-flame_2-2-Rの画像調整用の関数
def insert_ad_2_2_R(id, ad1, ad2):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(1.54), Cm(13.56), Cm(11.01))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(13.19), Cm(13.56), Cm(11.01))

# layout-flame_2-4-4-Lの画像調整用の関数
def insert_ad_2_4_4_L(id, ad1, ad2, ad3):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(3.09), Cm(1.54), Cm(13.56), Cm(11.01))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(3.09), Cm(13.24), Cm(13.56), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(3.09), Cm(19.04), Cm(13.56), Cm(5.12))

# layout-flame_2-4-4-Rの画像調整用の関数
def insert_ad_2_4_4_R(id, ad1, ad2, ad3):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(1.45), Cm(13.56), Cm(11.01))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(13.06), Cm(13.56), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(18.86), Cm(13.56), Cm(5.12))

# layout-flame_2-4-8-Lの画像調整用の関数
def insert_ad_2_4_8_L(id, ad1, ad2, ad3, ad4):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(3.09), Cm(1.54), Cm(13.56), Cm(11.01))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(3.09), Cm(13.24), Cm(13.56), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(3.09), Cm(19.04), Cm(6.63), Cm(5.12))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(9.87), Cm(19.04), Cm(6.63), Cm(5.12))

# layout-flame_2-4-8-Lの画像調整用の関数
def insert_ad_2_4_8_R(id, ad1, ad2, ad3, ad4):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.75), Cm(1.63), Cm(13.56), Cm(11.01))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.75), Cm(13.24), Cm(13.56), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(1.75), Cm(19.12), Cm(6.63), Cm(5.12))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(8.60), Cm(19.12), Cm(6.63), Cm(5.12))

# layout-flame_3-4-Lの画像調整用の関数
def insert_ad_3_4_L(id, ad1, ad2):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(1.54), Cm(13.56), Cm(16.84))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(19.04), Cm(13.56), Cm(5.12))

# layout-flame_3-4-Rの画像調整用の関数
def insert_ad_3_4_R(id, ad1, ad2):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(1.37), Cm(13.56), Cm(16.84))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(19.04), Cm(13.56), Cm(5.12))

# layout-flame_3-8-Lの画像調整用の関数
def insert_ad_3_8_L(id, ad1, ad2, ad3):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(1.54), Cm(13.56), Cm(16.84))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(19.04), Cm(6.63), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(9.87), Cm(19.04), Cm(6.63), Cm(5.12))

# layout-flame_3-8-Rの画像調整用の関数
def insert_ad_3_8_R(id, ad1, ad2, ad3):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(1.37), Cm(13.56), Cm(16.84))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(19.04), Cm(6.63), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(8.60), Cm(19.04), Cm(6.63), Cm(5.12))

# layout-flame_3(4)-4-Lの画像調整用の関数
def insert_ad_34_4_L(id, ad1, ad2, ad3, ad4):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(1.54), Cm(13.56), Cm(5.52))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(7.27), Cm(13.56), Cm(5.52))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(12.85), Cm(13.56), Cm(5.52))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(19.04), Cm(13.56), Cm(5.12))

# layout-flame_3(4)-4-Rの画像調整用の関数
def insert_ad_34_4_R(id, ad1, ad2, ad3, ad4):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(1.54), Cm(13.56), Cm(5.52))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(7.27), Cm(13.56), Cm(5.52))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(12.85), Cm(13.56), Cm(5.52))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(19.04), Cm(13.56), Cm(5.12))

# layout-flame_3(4)-8-Lの画像調整用の関数
def insert_ad_34_8_L(id, ad1, ad2, ad3, ad4, ad5):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(1.54), Cm(13.56), Cm(5.52))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(7.27), Cm(13.56), Cm(5.52))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(12.85), Cm(13.56), Cm(5.52))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(19.04), Cm(6.63), Cm(5.12))
    file_path = ad5
    slide.shapes.add_picture(file_path, Cm(9.87), Cm(19.04), Cm(6.63), Cm(5.12))

# layout-flame_3(4)-8-Rの画像調整用の関数
def insert_ad_34_8_R(id, ad1, ad2, ad3, ad4, ad5):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(1.54), Cm(13.56), Cm(5.52))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(7.27), Cm(13.56), Cm(5.52))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(12.85), Cm(13.56), Cm(5.52))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(1.66), Cm(19.04), Cm(6.63), Cm(5.12))
    file_path = ad5
    slide.shapes.add_picture(file_path, Cm(8.60), Cm(19.04), Cm(6.63), Cm(5.12))

# layout-flame_3(8)-4-Lの画像調整用の関数
def insert_ad_38_4_L(id, ad1, ad2, ad3, ad4, ad5, ad6, ad7):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(1.54), Cm(6.63), Cm(5.58))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(7.21), Cm(6.63), Cm(5.58))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(12.90), Cm(6.63), Cm(5.48))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(9.91), Cm(1.54), Cm(6.63), Cm(5.58))
    file_path = ad5
    slide.shapes.add_picture(file_path, Cm(9.91), Cm(7.21), Cm(6.63), Cm(5.58))
    file_path = ad6
    slide.shapes.add_picture(file_path, Cm(9.91), Cm(12.90), Cm(6.63), Cm(5.48))
    file_path = ad7
    slide.shapes.add_picture(file_path, Cm(2.97), Cm(19.04), Cm(13.56), Cm(5.12))

# layout-flame_3(8)-8-Rの画像調整用の関数
def insert_ad_38_8_R(id, ad1, ad2, ad3, ad4, ad5, ad6, ad7, ad8):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(1.54), Cm(6.63), Cm(5.58))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(7.22), Cm(6.63), Cm(5.58))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(12.85), Cm(6.63), Cm(5.58))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(1.68), Cm(18.99), Cm(6.63), Cm(5.12))
    file_path = ad5
    slide.shapes.add_picture(file_path, Cm(8.66), Cm(1.54), Cm(6.63), Cm(5.58))
    file_path = ad6
    slide.shapes.add_picture(file_path, Cm(8.66), Cm(7.22), Cm(6.63), Cm(5.58))
    file_path = ad7
    slide.shapes.add_picture(file_path, Cm(8.66), Cm(12.85), Cm(6.63), Cm(5.58))
    file_path = ad8
    slide.shapes.add_picture(file_path, Cm(8.66), Cm(18.99), Cm(6.63), Cm(5.12))

# layout-flame_4-4-4-4-Lの画像調整用の関数
def insert_ad_4_4_4_4_L(id, ad1, ad2, ad3, ad4):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(3.29), Cm(1.74), Cm(13.56), Cm(5.12))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(3.29), Cm(7.51), Cm(13.56), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(3.29), Cm(13.37), Cm(13.56), Cm(5.12))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(3.29), Cm(19.08), Cm(13.56), Cm(5.12))

# layout-flame_4-4-4-4-Rの画像調整用の関数
def insert_ad_4_4_4_4_R(id, ad1, ad2, ad3, ad4):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.31), Cm(1.58), Cm(13.56), Cm(5.12))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.31), Cm(7.38), Cm(13.56), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(1.31), Cm(13.18), Cm(13.56), Cm(5.12))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(1.31), Cm(18.98), Cm(13.56), Cm(5.12))

# layout-flame_4-4-4-8-Lの画像調整用の関数
def insert_ad_4_4_4_8_L(id, ad1, ad2, ad3, ad4, ad5):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(3.29), Cm(1.74), Cm(13.56), Cm(5.12))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(3.29), Cm(7.51), Cm(13.56), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(3.29), Cm(13.37), Cm(13.56), Cm(5.12))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(3.29), Cm(19.14), Cm(6.63), Cm(5.12))
    file_path = ad5
    slide.shapes.add_picture(file_path, Cm(10.15), Cm(19.14), Cm(6.63), Cm(5.12))

# layout-flame_4-4-4-8-Rの画像調整用の関数
def insert_ad_4_4_4_8_R(id, ad1, ad2, ad3, ad4, ad5):
    slide = pt.slides[id]  # layout-flame_something
    # 画像の挿入
    file_path = ad1
    slide.shapes.add_picture(file_path, Cm(1.31), Cm(1.58), Cm(13.56), Cm(5.12))
    file_path = ad2
    slide.shapes.add_picture(file_path, Cm(1.31), Cm(7.38), Cm(13.56), Cm(5.12))
    file_path = ad3
    slide.shapes.add_picture(file_path, Cm(1.31), Cm(13.18), Cm(13.56), Cm(5.12))
    file_path = ad4
    slide.shapes.add_picture(file_path, Cm(1.31), Cm(18.98), Cm(6.63), Cm(5.12))
    file_path = ad5
    slide.shapes.add_picture(file_path, Cm(8.20), Cm(18.98), Cm(6.63), Cm(5.12))

# ------------------------------------------------------------------------------------------------
# 関数で使う変数
ad_id_1 = 0
ad_id_2 = 0
ad_id_3 = 0
ad_id_4 = 0
ad_id_8 = 0
ads_1 = glob.glob('/Users/kosei/Desktop/transparented/*.1_*.png')  # xxx.1_xxx.pngという名前のファイルをさがす
ads_2 = glob.glob('/Users/kosei/Desktop/transparented/*.2_*.png') 
ads_3 = glob.glob('/Users/kosei/Desktop/transparented/*.3_*.png') 
ads_4 = glob.glob('/Users/kosei/Desktop/transparented/*.4_*.png') 
ads_8 = glob.glob('/Users/kosei/Desktop/transparented/*.8_*.png') 

# 新規ページを追加 → テキストボックス追加 → 広告追加
def greed_1_L(id):
    layout_flame = "layout-flame_1-L"
    # 該当するレイアウト枠を挿入した新規ページを追加
    add_page(layout_flame)
    ad = ads_1[ad_id_1]  # 添字により任意の写真を指定
    insert_textbox_1_L(id)
    insert_ad_1_L(id, ad)

def greed_1_R(id):
    layout_flame = "layout-flame_1-R"
    add_page(layout_flame)
    ad = ads_1[ad_id_1]
    insert_textbox_1_R(id)
    insert_ad_1_R(id, ad)

def greed_18_L(id):
    layout_flame = "layout-flame_1(8)-L"
    add_page(layout_flame)
    ad1 = ads_8[ad_id_8]
    ad2 = ads_8[ad_id_8]
    ad3 = ads_8[ad_id_8]
    ad4 = ads_8[ad_id_8]
    ad5 = ads_8[ad_id_8]
    ad6 = ads_8[ad_id_8]
    ad7 = ads_8[ad_id_8]
    ad8 = ads_8[ad_id_8]
    insert_textbox_1_L(id)
    insert_ad_18_L(id, ad1, ad2, ad3, ad4, ad5, ad6, ad7, ad8)

def greed_18_R(id):
    layout_flame = "layout-flame_1(8)-R"
    add_page(layout_flame)
    ad1 = ads_8[ad_id_8]
    ad2 = ads_8[ad_id_8]
    ad3 = ads_8[ad_id_8]
    ad4 = ads_8[ad_id_8]
    ad5 = ads_8[ad_id_8]
    ad6 = ads_8[ad_id_8]
    ad7 = ads_8[ad_id_8]
    ad8 = ads_8[ad_id_8]
    insert_textbox_1_R(id)
    insert_ad_18_R(id, ad1, ad2, ad3, ad4, ad5, ad6, ad7, ad8)

def greed_2_2_L(id):
    layout_flame = "layout-flame_2-2-L"
    add_page(layout_flame)
    ad1 = ads_2[ad_id_2]
    ad2 = ads_2[ad_id_2]
    insert_textbox_2_2_L(id)
    insert_ad_2_2_L(id, ad1, ad2)

def greed_2_2_R(id):
    layout_flame = "layout-flame_2-2-R"
    add_page(layout_flame)
    ad1 = ads_2[ad_id_2]
    ad2 = ads_2[ad_id_2]
    insert_textbox_2_2_R(id)
    insert_ad_2_2_R(id, ad1, ad2)

def greed_2_4_4_L(id):
    layout_flame = "layout-flame_2-4-4-L"
    add_page(layout_flame)
    ad1 = ads_2[ad_id_2]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    insert_textbox_2_4_L(id)
    insert_ad_2_4_4_L(id, ad1, ad2, ad3)

def greed_2_4_4_R(id):
    layout_flame = "layout-flame_2-4-4-R"
    add_page(layout_flame)
    ad1 = ads_2[ad_id_2]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    insert_textbox_2_4_R(id)
    insert_ad_2_4_4_R(id, ad1, ad2, ad3)

def greed_2_4_8_L(id):
    layout_flame = "layout-flame_2-4-8-L"
    add_page(layout_flame)
    ad1 = ads_2[ad_id_2]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_8[ad_id_8]
    ad4 = ads_8[ad_id_8] 
    insert_textbox_2_4_L(id)
    insert_ad_2_4_8_L(id, ad1, ad2, ad3, ad4)

def greed_2_4_8_R(id):
    layout_flame = "layout-flame_2-4-8-R"
    add_page(layout_flame)
    ad1 = ads_2[ad_id_2]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_8[ad_id_8]
    ad4 = ads_8[ad_id_8] 
    insert_textbox_2_4_R(id)
    insert_ad_2_4_8_R(id, ad1, ad2, ad3, ad4)

def greed_3_4_L(id):
    layout_flame = "layout-flame_3-4-L"
    add_page(layout_flame)
    ad1 = ads_3[ad_id_3]
    ad2 = ads_4[ad_id_4]
    insert_textbox_3_L(id)
    insert_ad_3_4_L(id, ad1, ad2)

def greed_3_4_R(id):
    layout_flame = "layout-flame_3-4-R"
    add_page(layout_flame)
    ad1 = ads_3[ad_id_3]
    ad2 = ads_4[ad_id_4]
    insert_textbox_3_R(id)
    insert_ad_3_4_R(id, ad1, ad2)

def greed_3_8_L(id):
    layout_flame = "layout-flame_3-8-L"
    add_page(layout_flame)
    ad1 = ads_3[ad_id_3]
    ad2 = ads_8[ad_id_8]
    ad3 = ads_8[ad_id_8] 
    insert_textbox_3_L(id)
    insert_ad_3_8_L(id, ad1, ad2, ad3)

def greed_3_8_R(id):
    layout_flame = "layout-flame_3-8-R"
    add_page(layout_flame)
    ad1 = ads_3[ad_id_3]
    ad2 = ads_8[ad_id_8]
    ad3 = ads_8[ad_id_8] 
    insert_textbox_3_R(id)
    insert_ad_3_8_R(id, ad1, ad2, ad3)

def greed_34_4_L(id):
    layout_flame = "layout-flame_3(4)-4-L"
    add_page(layout_flame)
    ad1 = ads_4[ad_id_4]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    ad4 = ads_4[ad_id_4]
    insert_textbox_3_L(id)
    insert_ad_34_4_L(id, ad1, ad2, ad3, ad4)

def greed_34_4_R(id):
    layout_flame = "layout-flame_3(4)-4-R"
    add_page(layout_flame)
    ad1 = ads_4[ad_id_4]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    ad4 = ads_4[ad_id_4]
    insert_textbox_3_R(id)
    insert_ad_34_4_R(id, ad1, ad2, ad3, ad4)

def greed_34_8_L(id):
    layout_flame = "layout-flame_3(4)-8-L"
    add_page(layout_flame)
    ad1 = ads_4[ad_id_4]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    ad4 = ads_8[ad_id_8]
    ad5 = ads_8[ad_id_8] 
    insert_textbox_3_L(id)
    insert_ad_34_8_L(id, ad1, ad2, ad3, ad4, ad5)

def greed_34_8_R(id):
    layout_flame = "layout-flame_3(4)-8-R"
    add_page(layout_flame)
    ad1 = ads_4[ad_id_4]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    ad4 = ads_8[ad_id_8]
    ad5 = ads_8[ad_id_8] 
    insert_textbox_3_R(id)
    insert_ad_34_8_R(id, ad1, ad2, ad3, ad4, ad5)

def greed_38_4_L(id):
    layout_flame = "layout-flame_3(8)-4-L"
    add_page(layout_flame)
    ad1 = ads_8[ad_id_8]
    ad2 = ads_8[ad_id_8]
    ad3 = ads_8[ad_id_8]
    ad4 = ads_8[ad_id_8]
    ad5 = ads_8[ad_id_8]
    ad6 = ads_8[ad_id_8]
    ad7 = ads_4[ad_id_4]
    insert_textbox_3_L(id)
    insert_ad_38_4_L(id, ad1, ad2, ad3, ad4, ad5, ad6, ad7)

def greed_38_8_R(id):
    layout_flame = "layout-flame_3(8)-8-R"
    add_page(layout_flame)
    ad1 = ads_8[ad_id_8]
    ad2 = ads_8[ad_id_8]
    ad3 = ads_8[ad_id_8]
    ad4 = ads_8[ad_id_8]
    ad5 = ads_8[ad_id_8]
    ad6 = ads_8[ad_id_8]
    ad7 = ads_8[ad_id_8]
    ad8 = ads_8[ad_id_8]
    insert_textbox_3_R(id)
    insert_ad_38_8_R(id, ad1, ad2, ad3, ad4, ad5, ad6, ad7, ad8)

def greed_4_4_4_4_L(id):
    layout_flame = "layout-flame_4-4-4-4-L"
    add_page(layout_flame)
    ad1 = ads_4[ad_id_4]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    ad4 = ads_4[ad_id_4]
    insert_textbox_4_L(id)
    insert_ad_4_4_4_4_L(id, ad1, ad2, ad3, ad4)

def greed_4_4_4_4_R(id):
    layout_flame = "layout-flame_4-4-4-4-R"
    add_page(layout_flame)
    ad1 = ads_4[ad_id_4]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    ad4 = ads_4[ad_id_4]
    insert_textbox_4_R(id)
    insert_ad_4_4_4_4_R(id, ad1, ad2, ad3, ad4)

def greed_4_4_4_8_L(id):
    layout_flame = "layout-flame_4-4-4-8-L"
    add_page(layout_flame)
    ad1 = ads_4[ad_id_4]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    ad4 = ads_8[ad_id_8]
    ad5 = ads_8[ad_id_8]
    insert_textbox_4_L(id)
    insert_ad_4_4_4_8_L(id, ad1, ad2, ad3, ad4, ad5)

def greed_4_4_4_8_R(id):
    layout_flame = "layout-flame_4-4-4-8-R"
    add_page(layout_flame)
    ad1 = ads_4[ad_id_4]
    ad2 = ads_4[ad_id_4]
    ad3 = ads_4[ad_id_4]
    ad4 = ads_8[ad_id_8]
    ad5 = ads_8[ad_id_8]
    insert_textbox_4_R(id)
    insert_ad_4_4_4_8_R(id, ad1, ad2, ad3, ad4, ad5)

# ------------------------------------------------------------------------------------------------
# Presentationクラスをインスタンス化
pt = Presentation()

# JIS B5 縦向き にスライドサイズを変更
pt.slide_height=Cm(25.7)
pt.slide_width=Cm(18.2)

# file name list
layout_flame_list = [
    "layout-flame_1-L",        # 1
    "layout-flame_1-R",        # 1
    "layout-flame_1(8)-L",     # 8
    "layout-flame_1(8)-R",     # 8
    "layout-flame_2-2-L",      # 2
    "layout-flame_2-2-R",      # 2
    "layout-flame_2-4-4-L",    # 3
    "layout-flame_2-4-4-R",    # 3
    "layout-flame_2-4-8-L",    # 4
    "layout-flame_2-4-8-R",    # 4
    "layout-flame_3-4-L",      # 2
    "layout-flame_3-4-R",      # 2
    "layout-flame_3-8-L",      # 3
    "layout-flame_3-8-R",      # 3
    "layout-flame_3(4)-4-L",   # 4
    "layout-flame_3(4)-4-R",   # 4
    "layout-flame_3(4)-8-L",   # 5
    "layout-flame_3(4)-8-R",   # 5
    "layout-flame_3(8)-4-L",   # 7
    "layout-flame_3(8)-8-R",   # 8
    "layout-flame_4-4-4-4-L",  # 4
    "layout-flame_4-4-4-4-R",  # 4
    "layout-flame_4-4-4-8-L",  # 5
    "layout-flame_4-4-4-8-R"   # 5
]

for i in range(10):
    num = random.randint(1, 4)
    print(num)
    if num == 1:
        greed_4_4_4_8_R(i)
    elif num == 2:
        greed_38_8_R(i)
    elif num == 3:
        greed_18_L(i)
    elif num == 4:
        greed_2_2_L(i)


# ファイルを任意の名前で保存（現在時刻をファイル名として保存するようにしている）
now = datetime.datetime.now()  # 現在時刻の取得
today = now.strftime('%Y年%m月%d日%H時%M分%S秒')  # 現在時刻を年月曜日で表示
save_name = '/Users/kosei/Desktop/output-files/' + today + '.pptx'  # 保存用のパワポのファイル名
pt.save(save_name)
print("ファイルの書き出しを完了しました")
